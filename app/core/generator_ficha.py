from __future__ import annotations

from collections.abc import Callable
from pathlib import Path
from typing import Final

from pptx import Presentation as PresentationFactory
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt

from app.core.reader import FichaEmployee, normalize_filename, parse_multiline_field

SLIDE_WIDTH: Final = Inches(13.333)
SLIDE_HEIGHT: Final = Inches(7.500)

VERDE_DETALHE: Final = "#84BD00"
VERDE_TITULO: Final = "#257226"
VERDE_PLACEHOLDER: Final = "#7CAF3D"
VERDE_PLACEHOLDER_PREENCHIMENTO: Final = "#F2F8EA"
BRANCO: Final = "#FFFFFF"
CINZA_META: Final = "#7A7A7A"
PRETO: Final = "#404040"

FONTE_TITULO: Final = "Calibri"
FONTE_CORPO: Final = "Calibri"

LEFT_ACCENT: Final = {"left": 0.0, "top": 0.0, "width": 0.183, "height": 2.283}
PHOTO_BOX: Final = {"left": 0.567, "top": 0.184, "width": 1.885, "height": 1.920}
TITLE_BOX: Final = {"left": 2.668, "top": 0.312, "width": 8.218, "height": 0.699}
META_BOX: Final = {"left": 2.657, "top": 0.873, "width": 5.200, "height": 0.942}
FOOTER_BAR: Final = {"left": 0.0, "top": 6.985, "width": 13.333, "height": 0.515}
BRAND_BLOCK: Final = {"left": 11.486, "top": 6.665, "width": 1.653, "height": 0.614}

FORMATION_TITLE: Final = {"left": 0.562, "top": 2.449, "width": 3.609, "height": 0.370}
FORMATION_RULE: Final = {"left": 0.580, "top": 2.819, "width": 5.000}
FORMATION_BODY: Final = {"left": 0.454, "top": 2.844, "width": 5.815, "height": 0.530}

SUMMARY_TITLE: Final = {"left": 6.395, "top": 2.436, "width": 3.609, "height": 0.370}
SUMMARY_RULE: Final = {"left": 6.413, "top": 2.806, "width": 6.099}
SUMMARY_BODY: Final = {"left": 6.417, "top": 2.844, "width": 6.722, "height": 1.100}

TRAJECTORY_TITLE: Final = {"left": 0.562, "top": 4.098, "width": 3.609, "height": 0.370}
TRAJECTORY_RULE: Final = {"left": 0.580, "top": 4.471, "width": 5.000}
TRAJECTORY_BODY: Final = {"left": 0.580, "top": 4.538, "width": 5.562, "height": 1.750}

PERFORMANCE_TITLE: Final = {"left": 6.402, "top": 5.407, "width": 3.609, "height": 0.370}
PERFORMANCE_RULE: Final = {"left": 6.413, "top": 5.745, "width": 6.099}
PERFORMANCE_BODY: Final = {"left": 6.417, "top": 5.867, "width": 5.562, "height": 0.858}


def create_presentation() -> Presentation:
    prs = PresentationFactory()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _rgb_from_hex(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _add_shape(
    slide: Slide,
    shape_type: MSO_AUTO_SHAPE_TYPE | MSO_SHAPE,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str | None = None,
    line_color: str | None = None,
    line_width_pt: float = 1.0,
):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb_from_hex(fill_color)

    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb_from_hex(line_color)
        shape.line.width = Pt(line_width_pt)
    return shape


def _set_run_font(
    run,
    *,
    size_pt: float,
    color: str,
    bold: bool = False,
    italic: bool = False,
    font_name: str = FONTE_CORPO,
) -> None:
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb_from_hex(color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name


def _new_textbox(
    slide: Slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
):
    text_box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = text_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = Inches(0.1)
    frame.margin_right = Inches(0.1)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    return text_box, frame


def _add_text(
    slide: Slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    color: str,
    size_pt: float,
    bold: bool = False,
    italic: bool = False,
    font_name: str = FONTE_CORPO,
    uppercase: bool = False,
) -> None:
    if text.strip() == "":
        return
    _textbox, frame = _new_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
    )
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text.upper() if uppercase else text
    _set_run_font(
        run,
        size_pt=size_pt,
        color=color,
        bold=bold,
        italic=italic,
        font_name=font_name,
    )


def _add_section_header(slide: Slide, *, title: str, title_box: dict[str, float], rule_box: dict[str, float]) -> None:
    _add_text(
        slide,
        text=title,
        left=title_box["left"],
        top=title_box["top"],
        width=title_box["width"],
        height=title_box["height"],
        color=VERDE_TITULO,
        size_pt=16,
        bold=True,
    )
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(rule_box["left"]),
        Inches(rule_box["top"]),
        Inches(rule_box["left"] + rule_box["width"]),
        Inches(rule_box["top"]),
    )
    connector.line.width = Pt(1.5)
    connector.line.color.rgb = _rgb_from_hex(VERDE_DETALHE)


def _add_photo_placeholder(slide: Slide) -> None:
    _add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=PHOTO_BOX["left"],
        top=PHOTO_BOX["top"],
        width=PHOTO_BOX["width"],
        height=PHOTO_BOX["height"],
        fill_color=VERDE_PLACEHOLDER_PREENCHIMENTO,
        line_color=VERDE_PLACEHOLDER,
        line_width_pt=2.0,
    )


def _clean(value: str | None) -> str:
    return "" if value is None else value.strip()


def _add_metadata_block(slide: Slide, employee: FichaEmployee) -> None:
    _textbox, frame = _new_textbox(
        slide,
        left=META_BOX["left"],
        top=META_BOX["top"],
        width=META_BOX["width"],
        height=META_BOX["height"],
    )
    lines = [
        _clean(employee.get("cargo")),
        f"Idade: {_clean(employee.get('idade'))} anos" if _clean(employee.get("idade")) else "",
        f"Tempo de Usiminas: {_clean(employee.get('antiguidade'))} anos"
        if _clean(employee.get("antiguidade"))
        else "",
    ]
    lines = [line for line in lines if line]
    if not lines:
        return

    for index, text in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = text
        _set_run_font(
            run,
            size_pt=16,
            color=CINZA_META,
            italic=True,
            font_name=FONTE_CORPO,
        )


def _add_body_text(slide: Slide, *, text: str, box: dict[str, float]) -> None:
    _add_text(
        slide,
        text=text,
        left=box["left"],
        top=box["top"],
        width=box["width"],
        height=box["height"],
        color=PRETO,
        size_pt=15,
        font_name=FONTE_CORPO,
    )


def _add_trajectory_block(slide: Slide, entries: list[str]) -> None:
    if not entries:
        return

    _textbox, frame = _new_textbox(
        slide,
        left=TRAJECTORY_BODY["left"],
        top=TRAJECTORY_BODY["top"],
        width=TRAJECTORY_BODY["width"],
        height=TRAJECTORY_BODY["height"],
    )
    frame.word_wrap = True

    for index, entry in enumerate(entries):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        cleaned = entry.strip()
        if cleaned.endswith(":"):
            run = paragraph.add_run()
            run.text = cleaned
            _set_run_font(
                run,
                size_pt=15,
                color=PRETO,
                bold=True,
                italic=True,
                font_name=FONTE_CORPO,
            )
            continue

        run = paragraph.add_run()
        run.text = cleaned
        _set_run_font(
            run,
            size_pt=15,
            color=PRETO,
            font_name=FONTE_CORPO,
        )


def _annual_notes(employee: FichaEmployee) -> list[tuple[str, str]]:
    values: list[tuple[str, str]] = []
    for field, label in (
        ("nota_2025", "2025"),
        ("nota_2024", "2024"),
        ("nota_2023", "2023"),
    ):
        value = _clean(employee.get(field))
        if value:
            values.append((label, value))
    return values


def _add_performance_block(slide: Slide, notes: list[tuple[str, str]]) -> None:
    if not notes:
        return

    _textbox, frame = _new_textbox(
        slide,
        left=PERFORMANCE_BODY["left"],
        top=PERFORMANCE_BODY["top"],
        width=PERFORMANCE_BODY["width"],
        height=PERFORMANCE_BODY["height"],
    )
    for index, (year, value) in enumerate(notes):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()

        bullet_run = paragraph.add_run()
        bullet_run.text = "• "
        _set_run_font(
            bullet_run,
            size_pt=15,
            color=PRETO,
            font_name=FONTE_CORPO,
        )

        year_run = paragraph.add_run()
        year_run.text = f"{year}:"
        _set_run_font(
            year_run,
            size_pt=15,
            color=PRETO,
            bold=True,
            font_name=FONTE_CORPO,
        )

        value_run = paragraph.add_run()
        value_run.text = f" {value}"
        _set_run_font(
            value_run,
            size_pt=15,
            color=PRETO,
            font_name=FONTE_CORPO,
        )


def build_slide(prs: Presentation, employee: FichaEmployee) -> Slide:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left=LEFT_ACCENT["left"],
        top=LEFT_ACCENT["top"],
        width=LEFT_ACCENT["width"],
        height=LEFT_ACCENT["height"],
        fill_color=VERDE_DETALHE,
    )
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left=FOOTER_BAR["left"],
        top=FOOTER_BAR["top"],
        width=FOOTER_BAR["width"],
        height=FOOTER_BAR["height"],
        fill_color=VERDE_TITULO,
    )
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left=BRAND_BLOCK["left"],
        top=BRAND_BLOCK["top"],
        width=BRAND_BLOCK["width"],
        height=BRAND_BLOCK["height"],
        fill_color=VERDE_DETALHE,
    )
    _add_text(
        slide,
        text="USIMINAS",
        left=11.630,
        top=6.786,
        width=1.300,
        height=0.260,
        color=BRANCO,
        size_pt=18,
        bold=True,
        font_name=FONTE_TITULO,
        uppercase=True,
    )

    _add_photo_placeholder(slide)

    _add_text(
        slide,
        text=_clean(employee.get("nome")),
        left=TITLE_BOX["left"],
        top=TITLE_BOX["top"],
        width=TITLE_BOX["width"],
        height=TITLE_BOX["height"],
        color=VERDE_TITULO,
        size_pt=36,
        bold=True,
        font_name=FONTE_TITULO,
        uppercase=True,
    )
    _add_metadata_block(slide, employee)

    formacao = _clean(employee.get("formacao"))
    if formacao:
        _add_section_header(slide, title="Formação", title_box=FORMATION_TITLE, rule_box=FORMATION_RULE)
        _add_body_text(slide, text=formacao, box=FORMATION_BODY)

    resumo = _clean(employee.get("resumo_perfil"))
    if resumo:
        _add_section_header(slide, title="Resumo", title_box=SUMMARY_TITLE, rule_box=SUMMARY_RULE)
        _add_body_text(slide, text=resumo, box=SUMMARY_BODY)

    trajetoria_items = parse_multiline_field(_clean(employee.get("trajetoria")))
    if trajetoria_items:
        _add_section_header(
            slide,
            title="Trajetória Profissional",
            title_box=TRAJECTORY_TITLE,
            rule_box=TRAJECTORY_RULE,
        )
        _add_trajectory_block(slide, trajetoria_items)

    notes = _annual_notes(employee)
    if notes:
        _add_section_header(
            slide,
            title="Performance e Potencial",
            title_box=PERFORMANCE_TITLE,
            rule_box=PERFORMANCE_RULE,
        )
        _add_performance_block(slide, notes)

    return slide


def _send_callback(callback: Callable[[dict], None] | None, payload: dict) -> None:
    if callback is not None:
        callback(payload)


def generate_ficha_pptx(
    employee: FichaEmployee,
    output_dir: str,
    *,
    callback: Callable[[dict], None] | None = None,
) -> str:
    output_root = Path(output_dir) / "fichas"
    output_root.mkdir(parents=True, exist_ok=True)
    name = _clean(employee.get("nome")) or "Colaborador"
    try:
        _send_callback(
            callback,
            {"type": "log", "message": f"Gerando ficha: {name}", "level": "success"},
        )
        prs = create_presentation()
        build_slide(prs, employee)
        file_stem = normalize_filename(name) or "Colaborador"
        output_path = output_root / f"{file_stem}.pptx"
        prs.save(str(output_path))
        _send_callback(
            callback,
            {"type": "progress", "current": 1, "total": 1, "name": name},
        )
        return str(output_path)
    except Exception as exc:
        _send_callback(
            callback,
            {
                "type": "log",
                "message": f"Erro ao gerar ficha para {name}: {exc}",
                "level": "error",
            },
        )
        raise
