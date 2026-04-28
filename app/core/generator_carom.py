from __future__ import annotations

import math
import re
from collections.abc import Callable
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from time import sleep
from typing import Any, Final, TypedDict

from pptx import Presentation as PresentationFactory
from pptx.slide import Slide

from app.core.carom_templates import CaromTemplate, get_carom_preset
from app.core.pptx_template_utils import (
    clear_text,
    clone_slide,
    replace_text,
    resolve_shape_path,
    reset_picture_to_circular_placeholder,
)
from app.core.reader import (
    CaromEmployee,
    resolve_carom_display_score_potential,
    validate_carom_employee_for_preset,
)

PROJETO_TRAINEE_BODY_TEXT: Final = "insira projeto trainee aqui"
CAROM_OUTPUT_TIMESTAMP_FORMAT: Final = "%d%m%Y_%H%M%S"
TALENT_REVIEW_STATIC_LINES: Final[tuple[str, str, str, str]] = (
    "Sucessor Imediato",
    "NomeCadeira",
    "Em desenvolvimento",
    "NomeCadeira",
)
TALENT_REVIEW_LINE_COUNT: Final = 6


class CaromConfig(TypedDict):
    preset_id: str
    titulo: str
    file_basename: str


def compute_projected_slide_count(selected_count: int, capacity: int) -> int:
    if selected_count <= 0:
        return 0
    return math.ceil(selected_count / max(capacity, 1))


def compute_current_slide_status(selected_count: int, capacity: int) -> str:
    safe_capacity = max(capacity, 1)
    if selected_count <= 0:
        return f"Faltam {safe_capacity} pessoas para completar o slide atual"
    position_in_current_slide = selected_count % safe_capacity
    if position_in_current_slide == 0:
        return "Slide atual completo"
    remaining = safe_capacity - position_in_current_slide
    return f"Faltam {remaining} pessoas para completar o slide atual"


def _clean(value: object) -> str:
    return "" if value is None else str(value).strip()


def _compose_non_empty(parts: list[str], separator: str) -> str:
    return separator.join(part for part in parts if part)


def _employee_name(employee: CaromEmployee) -> str:
    return _clean(employee.get("nome")) or "Sem Nome"


def _employee_age(employee: CaromEmployee) -> str:
    return _clean(employee.get("idade"))


def _employee_role(employee: CaromEmployee) -> str:
    return _clean(employee.get("cargo"))


def _employee_formation(employee: CaromEmployee) -> str:
    return _clean(employee.get("formacao"))


def _employee_ceo3(employee: CaromEmployee) -> str:
    return _clean(employee.get("ceo3"))


def _employee_ceo4(employee: CaromEmployee) -> str:
    return _clean(employee.get("ceo4"))


def _employee_score(employee: CaromEmployee) -> str:
    return _clean(resolve_carom_display_score_potential(employee))


def _clean_single_line(value: object) -> str:
    if value is None:
        return ""
    normalized = str(value).replace("\xa0", " ")
    return re.sub(r"\s+", " ", normalized).strip()


def _safe_file_basename(value: str, fallback: str) -> str:
    cleaned = _clean(value).replace("\xa0", " ")
    cleaned = re.sub(r"\s+", "_", cleaned)
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", cleaned)
    cleaned = re.sub(r"[._]+", "_", cleaned).strip("_")
    if cleaned == "":
        return fallback
    reserved_names = {
        "CON",
        "PRN",
        "AUX",
        "NUL",
        *(f"COM{index}" for index in range(1, 10)),
        *(f"LPT{index}" for index in range(1, 10)),
    }
    if cleaned.upper() in reserved_names:
        return f"_{cleaned}"
    return cleaned


def build_carom_output_filename(
    preset: CaromTemplate,
    generated_at: datetime | None = None,
    file_basename: str = "",
) -> str:
    timestamp = (generated_at or datetime.now()).strftime(CAROM_OUTPUT_TIMESTAMP_FORMAT)
    basename = _safe_file_basename(file_basename, f"Carometro{preset.output_type}")
    return f"{basename}_{timestamp}.pptx"


def _build_unique_carom_output_path(
    output_root: Path,
    preset: CaromTemplate,
    file_basename: str,
) -> Path:
    generated_at = datetime.now()
    output_path = output_root / build_carom_output_filename(
        preset,
        generated_at,
        file_basename,
    )
    while output_path.exists():
        current_second = generated_at.strftime(CAROM_OUTPUT_TIMESTAMP_FORMAT)
        while datetime.now().strftime(CAROM_OUTPUT_TIMESTAMP_FORMAT) == current_second:
            sleep(0.05)
        generated_at = datetime.now()
        output_path = output_root / build_carom_output_filename(
            preset,
            generated_at,
            file_basename,
        )
    return output_path


def _mini_lines(employee: CaromEmployee) -> list[str]:
    return [
        _employee_name(employee),
        _compose_non_empty([_employee_role(employee), _employee_age(employee)], " | "),
        _employee_ceo3(employee),
    ]


def _big_lines(employee: CaromEmployee) -> list[str]:
    headline = _compose_non_empty(
        [
            _employee_name(employee),
            _compose_non_empty(
                [_employee_age(employee), _employee_score(employee)],
                " - ",
            ),
        ],
        " | ",
    )
    return [
        headline,
        _employee_formation(employee),
        _employee_role(employee),
        _employee_ceo3(employee),
    ]


def _trainee_identity_lines(employee: CaromEmployee) -> list[str]:
    headline = _compose_non_empty(
        [
            _employee_name(employee),
            _compose_non_empty(
                [_employee_age(employee), _employee_score(employee)],
                " - ",
            ),
        ],
        " | ",
    )
    return [
        headline,
        _employee_role(employee),
        _employee_formation(employee),
        _employee_ceo4(employee),
    ]


def _talent_review_lines(employee: CaromEmployee) -> list[str]:
    headline = _compose_non_empty(
        [
            _clean_single_line(employee.get("nome")) or "Sem Nome",
            _compose_non_empty(
                [
                    _clean_single_line(employee.get("idade")),
                    _clean_single_line(resolve_carom_display_score_potential(employee)),
                ],
                " - ",
            ),
        ],
        " | ",
    )
    return [
        headline,
        _clean_single_line(employee.get("cargo")),
        *TALENT_REVIEW_STATIC_LINES,
    ]


def _replace_picture_at_path(slide: Slide, picture_path: tuple[int, ...]) -> None:
    picture_shape = resolve_shape_path(slide, picture_path)
    reset_picture_to_circular_placeholder(slide, picture_shape)


def _set_title_if_editable(slide: Slide, preset: CaromTemplate, title: str) -> None:
    if not preset.editable_title or preset.title_path is None:
        return
    replace_text(resolve_shape_path(slide, preset.title_path), [title])
    if preset.subtitle_path is not None:
        clear_text(resolve_shape_path(slide, preset.subtitle_path))


def _render_slot(
    slide: Slide,
    preset: CaromTemplate,
    slot: dict[str, tuple[int, ...]],
    employee: CaromEmployee | None,
) -> None:
    if preset.id == "mini":
        _render_basic_slot(slide, slot, employee, _mini_lines)
        return
    if preset.id == "big":
        _render_basic_slot(slide, slot, employee, _big_lines)
        return
    if preset.id == "projeto_trainee":
        _render_trainee_slot(slide, slot, employee)
        return
    if preset.id == "talent_review":
        _render_talent_review_slot(slide, slot, employee)
        return
    raise ValueError(f"Preset de carometro desconhecido: {preset.id}")


def _render_basic_slot(
    slide: Slide,
    slot: dict[str, tuple[int, ...]],
    employee: CaromEmployee | None,
    formatter: Callable[[CaromEmployee], list[str]],
) -> None:
    if employee is None:
        clear_text(resolve_shape_path(slide, slot["text"]))
        _replace_picture_at_path(slide, slot["picture"])
        return
    replace_text(resolve_shape_path(slide, slot["text"]), formatter(employee))
    _replace_picture_at_path(slide, slot["picture"])


def _render_trainee_slot(
    slide: Slide,
    slot: dict[str, tuple[int, ...]],
    employee: CaromEmployee | None,
) -> None:
    if employee is None:
        clear_text(resolve_shape_path(slide, slot["identity"]))
        clear_text(resolve_shape_path(slide, slot["body"]))
        _replace_picture_at_path(slide, slot["picture"])
        return
    replace_text(
        resolve_shape_path(slide, slot["identity"]), _trainee_identity_lines(employee)
    )
    replace_text(resolve_shape_path(slide, slot["body"]), [PROJETO_TRAINEE_BODY_TEXT])
    _replace_picture_at_path(slide, slot["picture"])


def _render_talent_review_slot(
    slide: Slide,
    slot: dict[str, tuple[int, ...]],
    employee: CaromEmployee | None,
) -> None:
    text_shape = _resolve_talent_review_text_shape(slide, slot["text"])
    if employee is None:
        _replace_talent_review_text(text_shape, [""] * TALENT_REVIEW_LINE_COUNT)
        _replace_picture_at_path(slide, slot["picture"])
        return
    _replace_talent_review_text(text_shape, _talent_review_lines(employee))
    _replace_picture_at_path(slide, slot["picture"])


def _replace_talent_review_text(shape: Any, lines: list[str]) -> None:
    frame = shape.text_frame
    styles = _capture_talent_review_styles(frame)
    _ensure_talent_review_paragraph_count(frame)

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[index]
        _replace_paragraph_with_style(paragraph, line, styles[index])


def _ensure_talent_review_paragraph_count(frame: Any) -> None:
    while len(frame.paragraphs) < TALENT_REVIEW_LINE_COUNT:
        frame.add_paragraph()

    while len(frame.paragraphs) > TALENT_REVIEW_LINE_COUNT:
        paragraph = frame.paragraphs[-1]
        paragraph._p.getparent().remove(paragraph._p)


def _capture_talent_review_styles(frame: Any) -> list[dict[str, Any]]:
    paragraphs = list(frame.paragraphs)
    styles = [
        _capture_paragraph_style(paragraphs[0]) if len(paragraphs) > 0 else {},
        _capture_paragraph_style(paragraphs[1]) if len(paragraphs) > 1 else {},
    ]
    styles.extend(_capture_static_talent_review_styles(paragraphs))
    while len(styles) < TALENT_REVIEW_LINE_COUNT:
        styles.append(styles[-1] if styles else {})
    return styles[:TALENT_REVIEW_LINE_COUNT]


def _capture_static_talent_review_styles(paragraphs: list[Any]) -> list[dict[str, Any]]:
    style_by_line: dict[int, dict[str, Any]] = {}
    nome_cadeira_indexes = [3, 5]

    for paragraph in paragraphs:
        for run in paragraph.runs:
            text = run.text.strip()
            if text == "Sucessor Imediato":
                style_by_line.setdefault(2, _capture_paragraph_style(paragraph, run))
            elif text == "NomeCadeira":
                next_index = nome_cadeira_indexes.pop(0) if nome_cadeira_indexes else 5
                style_by_line.setdefault(
                    next_index, _capture_paragraph_style(paragraph, run)
                )
            elif text == "Em desenvolvimento":
                style_by_line.setdefault(4, _capture_paragraph_style(paragraph, run))

    return [
        style_by_line.get(2) or _fallback_style(paragraphs, 2),
        style_by_line.get(3) or _fallback_style(paragraphs, 3),
        style_by_line.get(4) or _fallback_style(paragraphs, 4),
        style_by_line.get(5) or _fallback_style(paragraphs, 5),
    ]


def _fallback_style(paragraphs: list[Any], index: int) -> dict[str, Any]:
    if index < len(paragraphs):
        return _capture_paragraph_style(paragraphs[index])
    if paragraphs:
        return _capture_paragraph_style(paragraphs[-1])
    return {}


def _capture_paragraph_style(paragraph: Any, run: Any | None = None) -> dict[str, Any]:
    style: dict[str, Any] = {}
    if paragraph._p.pPr is not None:
        style["paragraph_properties"] = deepcopy(paragraph._p.pPr)

    source_run = run
    if source_run is None and paragraph.runs:
        source_run = paragraph.runs[0]
    if source_run is not None and source_run._r.rPr is not None:
        style["run_properties"] = deepcopy(source_run._r.rPr)
    return style


def _replace_paragraph_with_style(
    paragraph: Any, text: str, style: dict[str, Any]
) -> None:
    if paragraph._p.pPr is not None:
        paragraph._p.remove(paragraph._p.pPr)
    if style.get("paragraph_properties") is not None:
        paragraph._p.insert(0, deepcopy(style["paragraph_properties"]))

    paragraph_properties = paragraph._p.pPr
    for child in list(paragraph._p):
        if child is not paragraph_properties:
            paragraph._p.remove(child)
    if text == "":
        return

    run = paragraph.add_run()
    run.text = text
    if style.get("run_properties") is not None:
        if run._r.rPr is not None:
            run._r.remove(run._r.rPr)
        run._r.insert(0, deepcopy(style["run_properties"]))


def _resolve_talent_review_text_shape(slide: Slide, text_path: tuple[int, ...]):
    shape = resolve_shape_path(slide, text_path)
    if hasattr(shape, "text_frame"):
        return shape
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            if hasattr(child, "text_frame"):
                return child
    raise ValueError("Nao foi possivel localizar a caixa de texto do Talent Review.")


def _send_callback(callback: Callable[[dict], None] | None, payload: dict) -> None:
    if callback is not None:
        callback(payload)


def _validate_employees_for_preset(
    employees: list[CaromEmployee], preset_id: str
) -> None:
    for employee in employees:
        missing = validate_carom_employee_for_preset(employee, preset_id)
        if missing:
            joined = ", ".join(missing)
            name = _employee_name(employee)
            raise ValueError(
                f"O colaborador '{name}' nao possui os campos obrigatorios "
                f"para o template escolhido: {joined}."
            )


def generate_carom_pptx(
    employees: list[CaromEmployee],
    output_dir: str,
    config: CaromConfig,
    callback: Callable[[dict], None] | None = None,
) -> list[str]:
    if not employees:
        return []

    preset = get_carom_preset(config["preset_id"])
    _validate_employees_for_preset(employees, preset.id)
    title = _clean(config.get("titulo", "")) or preset.default_title
    output_root = Path(output_dir) / "carometros"
    output_root.mkdir(parents=True, exist_ok=True)

    prs = PresentationFactory(str(preset.template_path))
    template_slide = prs.slides[0]
    capacity = preset.capacity
    total_slides = compute_projected_slide_count(len(employees), capacity)
    slides: list[Slide] = [template_slide]
    for _ in range(total_slides - 1):
        slides.append(clone_slide(prs, template_slide))

    for slide_index, slide in enumerate(slides):
        start = slide_index * capacity
        end = start + capacity
        batch = employees[start:end]
        _set_title_if_editable(slide, preset, title)
        for slot_index, slot in enumerate(preset.slots):
            employee = batch[slot_index] if slot_index < len(batch) else None
            _render_slot(slide, preset, slot, employee)
        _send_callback(
            callback,
            {
                "type": "log",
                "message": f"Montando slide {slide_index + 1} de {total_slides}",
                "level": "info",
            },
        )
        _send_callback(
            callback,
            {
                "type": "progress",
                "current": slide_index + 1,
                "total": total_slides,
                "name": title,
            },
        )

    output_path = _build_unique_carom_output_path(
        output_root,
        preset,
        str(config.get("file_basename", "")),
    )
    prs.save(str(output_path))
    return [str(output_path)]
