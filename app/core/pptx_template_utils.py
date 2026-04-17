from __future__ import annotations

import struct
import zlib
from copy import deepcopy
from io import BytesIO
from typing import Any

from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.slide import Slide

_REL_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"))
PLACEHOLDER_RGB = RGBColor(232, 236, 241)


def clone_slide(prs, template_slide: Slide) -> Slide:
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in template_slide.shapes:
        copied = deepcopy(shape.element)
        _remap_relationships(template_slide.part, new_slide.part, copied)
        new_slide.shapes._spTree.insert_element_before(copied, "p:extLst")
    return new_slide


def _remap_relationships(source_part: Any, target_part: Any, element: Any) -> None:
    rel_map: dict[str, str] = {}
    for node in element.iter():
        for attr_name in _REL_ATTRS:
            old_rel_id = node.get(attr_name)
            if not old_rel_id:
                continue
            new_rel_id = rel_map.get(old_rel_id)
            if new_rel_id is None:
                relationship = source_part.rels[old_rel_id]
                new_rel_id = target_part.rels._add_relationship(
                    relationship.reltype,
                    relationship._target,
                    relationship.is_external,
                )
                rel_map[old_rel_id] = new_rel_id
            node.set(attr_name, new_rel_id)


def resolve_shape_path(container: Slide | Any, path: tuple[int, ...]):
    current = container
    for index in path:
        shapes = current.shapes
        current = shapes[index]
    return current


def replace_text(shape: Any, paragraphs: list[str]) -> None:
    frame = shape.text_frame
    while len(frame.paragraphs) < len(paragraphs):
        frame.add_paragraph()

    existing = list(frame.paragraphs)
    for index, paragraph in enumerate(existing):
        text = paragraphs[index] if index < len(paragraphs) else ""
        _replace_paragraph_text(paragraph, text)


def _replace_paragraph_text(paragraph: Any, text: str) -> None:
    style = _capture_run_style(paragraph)
    for run in list(paragraph.runs):
        paragraph._p.remove(run._r)
    if text == "":
        return
    run = paragraph.add_run()
    run.text = text
    _apply_run_style(run, style)


def _capture_run_style(paragraph: Any) -> dict[str, Any]:
    if paragraph.runs:
        font = paragraph.runs[0].font
        return {
            "size": font.size,
            "bold": font.bold,
            "italic": font.italic,
            "name": font.name,
            "underline": font.underline,
            "color": _capture_font_color(font),
        }
    return {}


def _capture_font_color(font: Any) -> RGBColor | None:
    try:
        color = font.color.rgb
    except Exception:
        return None
    return color


def _apply_run_style(run: Any, style: dict[str, Any]) -> None:
    font = run.font
    if style.get("size") is not None:
        font.size = style["size"]
    if style.get("bold") is not None:
        font.bold = style["bold"]
    if style.get("italic") is not None:
        font.italic = style["italic"]
    if style.get("name"):
        font.name = style["name"]
    if style.get("underline") is not None:
        font.underline = style["underline"]
    if style.get("color") is not None:
        font.color.rgb = style["color"]


def clear_text(shape: Any) -> None:
    replace_text(shape, [])


def replace_picture(slide: Slide, picture_shape: Any, image_bytes: bytes) -> Any:
    image_stream = BytesIO(image_bytes)
    new_picture = slide.shapes.add_picture(
        image_stream,
        picture_shape.left,
        picture_shape.top,
        picture_shape.width,
        picture_shape.height,
    )
    picture_shape._element.addprevious(new_picture._element)
    picture_shape._element.getparent().remove(picture_shape._element)
    return new_picture


def reset_picture_to_circular_placeholder(slide: Slide, picture_shape: Any) -> Any:
    diameter = min(picture_shape.width, picture_shape.height)
    picture_shape.left = picture_shape.left + (picture_shape.width - diameter) // 2
    picture_shape.top = picture_shape.top + (picture_shape.height - diameter) // 2
    picture_shape.width = diameter
    picture_shape.height = diameter

    image_stream = BytesIO(placeholder_picture_bytes())
    _image_part, rel_id = slide.part.get_or_add_image_part(image_stream)
    blip = picture_shape._element.blipFill.blip
    if blip is not None:
        blip.set(qn("r:embed"), rel_id)
        blip.attrib.pop(qn("r:link"), None)

    src_rect = picture_shape._element.blipFill.srcRect
    if src_rect is not None:
        src_rect.getparent().remove(src_rect)

    sp_pr = picture_shape._element.spPr
    if sp_pr.custGeom is not None:
        sp_pr._remove_custGeom()
    prst_geom = sp_pr.prstGeom if sp_pr.prstGeom is not None else sp_pr._add_prstGeom()
    prst_geom.set("prst", "ellipse")
    if prst_geom.avLst is None:
        prst_geom._add_avLst()
    return picture_shape


def placeholder_picture_bytes() -> bytes:
    return _solid_png_bytes(PLACEHOLDER_RGB[0], PLACEHOLDER_RGB[1], PLACEHOLDER_RGB[2])


def _solid_png_bytes(red: int, green: int, blue: int) -> bytes:
    raw_scanline = b"\x00" + bytes((red, green, blue, 255))
    compressed = zlib.compress(raw_scanline)
    chunks = [
        _png_chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0)),
        _png_chunk(b"IDAT", compressed),
        _png_chunk(b"IEND", b""),
    ]
    return b"\x89PNG\r\n\x1a\n" + b"".join(chunks)


def _png_chunk(chunk_type: bytes, data: bytes) -> bytes:
    checksum = zlib.crc32(chunk_type + data) & 0xFFFFFFFF
    return struct.pack(">I", len(data)) + chunk_type + data + struct.pack(">I", checksum)
