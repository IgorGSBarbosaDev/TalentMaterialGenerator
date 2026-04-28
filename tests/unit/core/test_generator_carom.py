from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from app.core import generator_carom
from app.core.carom_templates import CAROM_TEMPLATES, get_carom_preset
from app.core.pptx_template_utils import resolve_shape_path


def _employee(index: int) -> dict[str, str]:
    return {
        "matricula": str(100 + index),
        "nome": f"Colab {index}",
        "idade": str(20 + index),
        "cargo": "Analista",
        "formacao": "Engenharia",
        "resumo_perfil": "Resumo",
        "trajetoria": "Trajetoria",
        "foto": "",
        "area": "Operacao",
        "localizacao": "",
        "unidade_gestao": "",
        "nota_2025": "4 / AP",
        "avaliacao_2025": "4 / AP",
        "score_2025": "4",
        "potencial_2025": "AP",
        "ceo3": f"CEO3 {index}",
        "ceo4": f"CEO4 {index}",
    }


def _all_slide_text(slide) -> str:
    values: list[str] = []

    def _collect(shape) -> None:
        if hasattr(shape, "text"):
            values.append(shape.text)
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                _collect(child)

    for shape in slide.shapes:
        _collect(shape)
    return "\n".join(values)


def _assert_picture_slots_are_circular_placeholders(
    file_path: str, preset_id: str
) -> None:
    prs = Presentation(file_path)
    preset = get_carom_preset(preset_id)
    for slot in preset.slots:
        shape = resolve_shape_path(prs.slides[0], slot["picture"])
        prst_geom = shape._element.spPr.prstGeom
        src_rect = shape._element.blipFill.srcRect
        assert shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        assert prst_geom is not None
        assert prst_geom.get("prst") == "ellipse"
        assert shape.width == shape.height
        assert src_rect is None or not src_rect.attrib


def _talent_review_slot_paragraphs(file_path: str, slot_index: int) -> list[str]:
    prs = Presentation(file_path)
    text_shape = _talent_review_slot_text_shape(prs, slot_index)
    return [paragraph.text for paragraph in text_shape.text_frame.paragraphs]


def _talent_review_slot_text_shape(prs: Presentation, slot_index: int):
    preset = get_carom_preset("talent_review")
    shape = resolve_shape_path(prs.slides[0], preset.slots[slot_index]["text"])
    if not hasattr(shape, "text_frame") and hasattr(shape, "shapes"):
        shape = next(child for child in shape.shapes if hasattr(child, "text_frame"))
    return shape


def _paragraph_run_properties(file_path: str, slot_index: int, paragraph_index: int):
    prs = Presentation(file_path)
    text_shape = _talent_review_slot_text_shape(prs, slot_index)
    paragraph = text_shape.text_frame.paragraphs[paragraph_index]
    assert len(paragraph.runs) == 1
    run_properties = paragraph.runs[0]._r.rPr
    assert run_properties is not None
    return run_properties


def _run_srgb_fill(run_properties) -> str | None:
    srgb = run_properties.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    return None if srgb is None else srgb.get("val")


def _run_scheme_fill(run_properties) -> str | None:
    scheme = run_properties.find(qn("a:solidFill")).find(qn("a:schemeClr"))
    return None if scheme is None else scheme.get("val")


def _run_highlight(run_properties) -> str | None:
    highlight = run_properties.find(qn("a:highlight"))
    if highlight is None:
        return None
    srgb = highlight.find(qn("a:srgbClr"))
    return None if srgb is None else srgb.get("val")


def _talent_review_expected_lines(employee: dict[str, str]) -> list[str]:
    return [
        f"{employee['nome']} | {employee['idade']} - {employee['nota_2025']}",
        employee["cargo"],
        "Sucessor Imediato",
        "NomeCadeira",
        "Em desenvolvimento",
        "NomeCadeira",
    ]


def test_get_carom_preset_maps_legacy_regular_to_mini_template() -> None:
    preset = generator_carom.get_carom_preset("regular")

    assert preset.id == "mini"
    assert preset.capacity == 18
    assert preset.template_path.name == "Carometro-mini.pptx"


def test_get_carom_preset_maps_legacy_large_to_big_template() -> None:
    preset = generator_carom.get_carom_preset("large")

    assert preset.id == "big"
    assert preset.capacity == 8
    assert preset.template_path.name == "Carometro-big.pptx"


def test_compute_projected_slide_count_uses_capacity() -> None:
    assert generator_carom.compute_projected_slide_count(0, 8) == 0
    assert generator_carom.compute_projected_slide_count(8, 8) == 1
    assert generator_carom.compute_projected_slide_count(23, 8) == 3


def test_compute_current_slide_status_reports_remaining_people() -> None:
    assert (
        generator_carom.compute_current_slide_status(3, 8)
        == "Faltam 5 pessoas para completar o slide atual"
    )
    assert generator_carom.compute_current_slide_status(8, 8) == "Slide atual completo"
    assert (
        generator_carom.compute_current_slide_status(9, 8)
        == "Faltam 7 pessoas para completar o slide atual"
    )


def test_build_carom_output_filename_uses_preset_type_and_local_timestamp() -> None:
    filename = generator_carom.build_carom_output_filename(
        get_carom_preset("mini"),
        datetime(2026, 4, 16, 14, 35, 22),
    )

    assert filename == "CarometroMini_16042026_143522.pptx"


def test_carom_output_types_match_required_filenames() -> None:
    generated_at = datetime(2026, 4, 16, 14, 35, 22)

    assert {
        preset_id: generator_carom.build_carom_output_filename(preset, generated_at)
        for preset_id, preset in CAROM_TEMPLATES.items()
    } == {
        "mini": "CarometroMini_16042026_143522.pptx",
        "big": "CarometroBig_16042026_143522.pptx",
        "projeto_trainee": "CarometroProjetoTrainee_16042026_143522.pptx",
        "talent_review": "CarometroTalentReview_16042026_143522.pptx",
    }


def test_build_carom_output_filename_changes_between_seconds() -> None:
    preset = get_carom_preset("big")

    first = generator_carom.build_carom_output_filename(
        preset,
        datetime(2026, 4, 16, 14, 35, 22),
    )
    second = generator_carom.build_carom_output_filename(
        preset,
        datetime(2026, 4, 16, 14, 35, 23),
    )

    assert first != second


def test_unique_carom_output_path_waits_for_next_second_when_file_exists(
    monkeypatch,
    tmp_path: Path,
) -> None:
    class FakeDateTime:
        values = iter(
            (
                datetime(2026, 4, 16, 14, 35, 22),
                datetime(2026, 4, 16, 14, 35, 22),
                datetime(2026, 4, 16, 14, 35, 23),
                datetime(2026, 4, 16, 14, 35, 23),
            )
        )

        @classmethod
        def now(cls):
            return next(cls.values)

    preset = get_carom_preset("big")
    existing = tmp_path / "CarometroBig_16042026_143522.pptx"
    existing.write_bytes(b"existing")
    monkeypatch.setattr(generator_carom, "datetime", FakeDateTime)
    monkeypatch.setattr(generator_carom, "sleep", lambda _seconds: None)

    output_path = generator_carom._build_unique_carom_output_path(
        tmp_path,
        preset,
        "",
    )

    assert output_path == tmp_path / "CarometroBig_16042026_143523.pptx"


def test_generate_carom_pptx_creates_single_file(tmp_path: Path) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(1), _employee(2)],
        str(tmp_path),
        {
            "preset_id": "big",
            "titulo": "Leadership Board",
            "file_basename": "Leadership_Board",
        },
    )

    assert len(files) == 1
    assert Path(files[0]).exists()


def test_generate_carom_pptx_uses_timestamped_configured_output_name(
    tmp_path: Path,
) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(1)],
        str(tmp_path),
        {
            "preset_id": "big",
            "titulo": "Leadership Board",
            "file_basename": "Leadership_Board",
        },
    )

    output_path = Path(files[0])
    assert output_path.parent == tmp_path / "carometros"
    assert re.fullmatch(r"Leadership_Board_\d{8}_\d{6}\.pptx", output_path.name)


def test_generate_carom_pptx_sanitizes_configured_output_name(
    tmp_path: Path,
) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(1)],
        str(tmp_path),
        {
            "preset_id": "big",
            "titulo": "Ignored",
            "file_basename": r"..\Talent/Review:Ciclo*2026?",
        },
    )

    output_path = Path(files[0])
    assert output_path.parent == tmp_path / "carometros"
    assert re.fullmatch(
        r"Talent_Review_Ciclo_2026_\d{8}_\d{6}\.pptx",
        output_path.name,
    )


def test_generate_carom_pptx_breaks_big_selection_into_multiple_slides(
    tmp_path: Path,
) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(index) for index in range(1, 10)],
        str(tmp_path),
        {
            "preset_id": "big",
            "titulo": "Leadership Board",
            "file_basename": "Leadership_Board",
        },
    )

    prs = Presentation(files[0])
    assert len(prs.slides) == 2


@pytest.mark.parametrize("preset_id", tuple(CAROM_TEMPLATES))
def test_generate_carom_pptx_uses_circular_photo_placeholders_for_every_preset(
    tmp_path: Path,
    preset_id: str,
) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(index) for index in range(1, 3)],
        str(tmp_path),
        {"preset_id": preset_id, "titulo": "Carometro", "file_basename": "Carometro"},
    )

    _assert_picture_slots_are_circular_placeholders(files[0], preset_id)


def test_generate_carom_pptx_ignores_valid_foto_path_for_manual_placeholders(
    tmp_path: Path,
) -> None:
    employee = _employee(1)
    employee["foto"] = str(Path("tests/fixtures/fotos/avatar_test.png").resolve())

    files = generator_carom.generate_carom_pptx(
        [employee],
        str(tmp_path),
        {
            "preset_id": "big",
            "titulo": "Leadership Board",
            "file_basename": "Leadership_Board",
        },
    )

    _assert_picture_slots_are_circular_placeholders(files[0], "big")


def test_generate_carom_pptx_uses_title_on_every_big_slide(tmp_path: Path) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(index) for index in range(1, 12)],
        str(tmp_path),
        {
            "preset_id": "big",
            "titulo": "Leadership Board",
            "file_basename": "Leadership_Board",
        },
    )

    prs = Presentation(files[0])
    slide_titles = [slide.shapes[17].text for slide in prs.slides]
    assert slide_titles == ["Leadership Board", "Leadership Board"]


def test_generate_carom_pptx_allows_legacy_regular_without_ceo_fields(
    tmp_path: Path,
) -> None:
    employee = _employee(1)
    employee["ceo3"] = ""
    employee["ceo4"] = ""
    employee["nota_2025"] = ""
    employee["avaliacao_2025"] = ""
    employee["score_2025"] = ""
    employee["potencial_2025"] = ""

    files = generator_carom.generate_carom_pptx(
        [employee],
        str(tmp_path),
        {"preset_id": "regular", "titulo": "Carometro", "file_basename": "Carometro"},
    )

    assert Path(files[0]).exists()


def test_generate_carom_pptx_rejects_big_when_required_ceo_field_is_missing(
    tmp_path: Path,
) -> None:
    employee = _employee(1)
    employee["ceo3"] = ""

    with pytest.raises(ValueError, match="ceo3"):
        generator_carom.generate_carom_pptx(
            [employee],
            str(tmp_path),
            {
                "preset_id": "big",
                "titulo": "Leadership Board",
                "file_basename": "Leadership_Board",
            },
        )


def test_generate_carom_pptx_clears_unused_big_slots_without_sample_text(
    tmp_path: Path,
) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(1)],
        str(tmp_path),
        {
            "preset_id": "big",
            "titulo": "Leadership Board",
            "file_basename": "Leadership_Board",
        },
    )

    prs = Presentation(files[0])
    text = _all_slide_text(prs.slides[0])

    assert "Ana Carolina Alves Melo Gouveia" not in text
    assert "Fernando Augusto Rodrigues Machado" not in text
    assert "Colab 1" in text


def test_generate_carom_pptx_uses_literal_body_text_for_projeto_trainee(
    tmp_path: Path,
) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(1), _employee(2)],
        str(tmp_path),
        {
            "preset_id": "projeto_trainee",
            "titulo": "Ignored",
            "file_basename": "Projeto_Trainee",
        },
    )

    prs = Presentation(files[0])
    slide = prs.slides[0]

    assert slide.shapes[9].text.splitlines()[0] == "insira projeto trainee aqui"
    assert slide.shapes[13].text.splitlines()[0] == "insira projeto trainee aqui"
    assert slide.shapes[3].text == "CEO3 - Programa GT & GP"


def test_generate_carom_pptx_keeps_talent_review_template_text_without_ceo_fields(
    tmp_path: Path,
) -> None:
    employee = _employee(1)
    employee["ceo3"] = ""
    employee["ceo4"] = ""

    files = generator_carom.generate_carom_pptx(
        [employee],
        str(tmp_path),
        {
            "preset_id": "talent_review",
            "titulo": "Ignored",
            "file_basename": "Talent_Review",
        },
    )

    prs = Presentation(files[0])
    slide = prs.slides[0]
    slide_text = _all_slide_text(slide)
    text_box = slide.shapes[5]
    paragraphs = [paragraph.text for paragraph in text_box.text_frame.paragraphs]

    assert paragraphs[0] == "Colab 1 | 21 - 4 / AP"
    assert paragraphs[1] == "Analista"
    assert paragraphs[2] == "Sucessor Imediato"
    assert paragraphs[3] == "NomeCadeira"
    assert paragraphs[4] == "Em desenvolvimento"
    assert paragraphs[5] == "NomeCadeira"
    assert "CEO3" not in slide_text
    assert "CEO4" not in slide_text


def test_generate_carom_pptx_rebuilds_every_talent_review_slot_with_static_successor_block(
    tmp_path: Path,
) -> None:
    employee = _employee(1)
    employee.update(
        {
            "nome": "ADEILSON ROBERTO MENDES",
            "idade": "56",
            "cargo": "COMMERCIAL EXECUTIVE MANAGER",
            "nota_2025": "3 / MN+",
            "avaliacao_2025": "3 / MN+",
            "score_2025": "3",
            "potencial_2025": "MN+",
            "ceo3": "Commercial Vice-Presidency - USI",
            "ceo4": "Commercial - USI",
        }
    )
    employees = [dict(employee) for _slot in get_carom_preset("talent_review").slots]

    files = generator_carom.generate_carom_pptx(
        employees,
        str(tmp_path),
        {
            "preset_id": "talent_review",
            "titulo": "Ignored",
            "file_basename": "Talent_Review",
        },
    )

    prs = Presentation(files[0])
    slide_text = _all_slide_text(prs.slides[0])
    expected_lines = _talent_review_expected_lines(employee)
    for slot_index in range(len(get_carom_preset("talent_review").slots)):
        paragraphs = _talent_review_slot_paragraphs(files[0], slot_index)
        assert paragraphs == expected_lines
        assert paragraphs[0] != ""
        assert paragraphs[1] != ""

    assert "\x0b" not in slide_text
    assert "Commercial Vice-Presidency - USI" not in slide_text
    assert "Commercial - USI" not in slide_text
    assert "CEO3" not in slide_text
    assert "CEO4" not in slide_text


def test_generate_carom_pptx_preserves_talent_review_static_run_formatting(
    tmp_path: Path,
) -> None:
    employee = _employee(1)
    employees = [dict(employee) for _slot in get_carom_preset("talent_review").slots]

    files = generator_carom.generate_carom_pptx(
        employees,
        str(tmp_path),
        {
            "preset_id": "talent_review",
            "titulo": "Ignored",
            "file_basename": "Talent_Review",
        },
    )

    for slot_index in range(len(get_carom_preset("talent_review").slots)):
        successor = _paragraph_run_properties(files[0], slot_index, 2)
        first_nomination = _paragraph_run_properties(files[0], slot_index, 3)
        development = _paragraph_run_properties(files[0], slot_index, 4)
        second_nomination = _paragraph_run_properties(files[0], slot_index, 5)

        assert _run_highlight(successor) == "84BD00"
        assert _run_scheme_fill(successor) == "bg1"
        assert _run_highlight(development) == "0000FF"
        assert _run_scheme_fill(development) == "bg1"

        assert _run_highlight(first_nomination) is None
        assert _run_srgb_fill(first_nomination) == "242424"
        assert _run_highlight(second_nomination) is None
        assert _run_srgb_fill(second_nomination) == "242424"


def test_generate_carom_pptx_normalizes_talent_review_dynamic_text(
    tmp_path: Path,
) -> None:
    employee = _employee(1)
    employee.update(
        {
            "nome": "\tADEILSON\r\nROBERTO\xa0\xa0MENDES ",
            "idade": " 56 ",
            "cargo": "\nCOMMERCIAL\t EXECUTIVE\x0bMANAGER  ",
            "nota_2025": "",
            "avaliacao_2025": "",
            "score_2025": " 3 ",
            "potencial_2025": " MN+ ",
        }
    )

    files = generator_carom.generate_carom_pptx(
        [employee],
        str(tmp_path),
        {
            "preset_id": "talent_review",
            "titulo": "Ignored",
            "file_basename": "Talent_Review",
        },
    )

    paragraphs = _talent_review_slot_paragraphs(files[0], 0)
    assert paragraphs == [
        "ADEILSON ROBERTO MENDES | 56 - 3 / MN+",
        "COMMERCIAL EXECUTIVE MANAGER",
        "Sucessor Imediato",
        "NomeCadeira",
        "Em desenvolvimento",
        "NomeCadeira",
    ]


def test_generate_carom_pptx_paginates_full_talent_review_capacity(
    tmp_path: Path,
) -> None:
    files = generator_carom.generate_carom_pptx(
        [_employee(index) for index in range(1, 14)],
        str(tmp_path),
        {
            "preset_id": "talent_review",
            "titulo": "Ignored",
            "file_basename": "Talent_Review",
        },
    )

    prs = Presentation(files[0])
    assert len(prs.slides) == 2
