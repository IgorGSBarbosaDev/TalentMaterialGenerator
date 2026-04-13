from __future__ import annotations

from pathlib import Path

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from app.core import generator_ficha
from app.core.reader import FichaEmployee, normalize_filename


def _employee(**overrides: str) -> FichaEmployee:
    base: FichaEmployee = {
        "matricula": "123",
        "nome": "Ana Martins",
        "idade": "30",
        "cargo": "Analista",
        "antiguidade": "5",
        "formacao": "Engenharia",
        "trajetoria": "Usiminas:; 2024-2025 - Coordenadora; Externo:; 2022-2024 - Analista",
        "resumo_perfil": "Resumo profissional",
        "nota_2025": "4 / PROM",
        "nota_2024": "3 / PROM",
        "nota_2023": "5 / PROM",
        "avaliacao_2025": "",
        "avaliacao_2024": "",
        "avaliacao_2023": "",
        "score_2025": "",
        "score_2024": "",
        "score_2023": "",
        "potencial_2025": "",
        "potencial_2024": "",
        "potencial_2023": "",
    }
    base.update(overrides)
    return base


def _shape_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            value = shape.text.strip()
            if value:
                texts.append(value)
    return texts


def _shape_in_inches(shape) -> tuple[float, float, float, float]:
    inch = Inches(1)
    return (
        round(shape.left / inch, 3),
        round(shape.top / inch, 3),
        round(shape.width / inch, 3),
        round(shape.height / inch, 3),
    )


def _find_auto_shape(slide, auto_shape_type, *, left: float, top: float, tol: float = 0.03):
    for shape in slide.shapes:
        try:
            shape_auto_type = shape.auto_shape_type
        except (AttributeError, ValueError):
            continue
        if shape_auto_type != auto_shape_type:
            continue
        shape_left, shape_top, _shape_width, _shape_height = _shape_in_inches(shape)
        if abs(shape_left - left) <= tol and abs(shape_top - top) <= tol:
            return shape
    return None


def test_build_slide_creates_slide() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(prs, _employee())

    assert slide is not None
    assert len(prs.slides) == 1


def test_build_slide_matches_reference_canvas_size() -> None:
    prs = generator_ficha.create_presentation()

    assert prs.slide_width == Inches(13.333)
    assert prs.slide_height == Inches(7.5)


def test_build_slide_omits_empty_optional_sections() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(
        prs,
        _employee(formacao="", trajetoria="", resumo_perfil="", nota_2025="", nota_2024="", nota_2023=""),
    )

    texts = [text.upper() for text in _shape_texts(slide)]
    assert "FORMAÇÃO" not in texts
    assert "TRAJETÓRIA PROFISSIONAL" not in texts
    assert "RESUMO" not in texts
    assert "PERFORMANCE E POTENCIAL" not in texts


def test_build_slide_uses_rounded_photo_placeholder_instead_of_oval() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(prs, _employee())

    oval_placeholder = _find_auto_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left=0.567,
        top=0.184,
    )
    rounded_placeholder = _find_auto_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=0.567,
        top=0.184,
    )

    assert oval_placeholder is not None
    assert rounded_placeholder is None


def test_build_slide_includes_reference_section_titles() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(prs, _employee())

    texts = _shape_texts(slide)

    assert "Formação" in texts
    assert "Resumo" in texts
    assert "Trajetória Profissional" in texts
    assert "Performance e Potencial" in texts


def test_build_slide_renders_annual_notes_from_structured_fallback() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(
        prs,
        _employee(
            nota_2025="",
            nota_2024="",
            nota_2023="",
            score_2025="4",
            potencial_2025="PROM",
            score_2024="",
            potencial_2024="AP",
            score_2023="5",
            potencial_2023="",
        ),
    )

    texts = _shape_texts(slide)

    assert any("2025: 4 / PROM" in text for text in texts)
    assert any("2024: AP" in text for text in texts)
    assert any("2023: 5" in text for text in texts)


def test_build_slide_ignores_invalid_annual_note_values() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(
        prs,
        _employee(
            nota_2025="",
            nota_2024="",
            nota_2023="",
            avaliacao_2025="#N/A",
            score_2025="#N/A",
            potencial_2025="#N/A",
            score_2024="",
            potencial_2024="",
            score_2023="",
            potencial_2023="",
        ),
    )

    texts = [text.upper() for text in _shape_texts(slide)]

    assert "PERFORMANCE E POTENCIAL" not in texts


def test_build_slide_contains_reference_geometry_landmarks() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(prs, _employee())

    left_accent = _find_auto_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left=0.0, top=0.0)
    footer = _find_auto_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left=0.0, top=6.985)
    brand = _find_auto_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left=11.486, top=6.665)
    placeholder = _find_auto_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left=0.567,
        top=0.184,
    )

    assert left_accent is not None
    assert _shape_in_inches(left_accent) == (0.0, 0.0, 0.183, 2.283)
    assert footer is not None
    assert _shape_in_inches(footer) == (0.0, 6.985, 13.333, 0.515)
    assert brand is not None
    assert _shape_in_inches(brand) == (11.486, 6.665, 1.653, 0.614)
    assert placeholder is not None
    assert _shape_in_inches(placeholder) == (0.567, 0.184, 1.885, 1.885)


def test_build_slide_uses_21pt_usiminas_label() -> None:
    prs = generator_ficha.create_presentation()
    slide = generator_ficha.build_slide(prs, _employee())

    usiminas_shape = next(
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip() == "USIMINAS"
    )
    run = usiminas_shape.text_frame.paragraphs[0].runs[0]

    assert round(run.font.size.pt, 1) == 21.0


def test_generate_ficha_pptx_creates_single_file(tmp_path: Path) -> None:
    created = generator_ficha.generate_ficha_pptx(_employee(nome="Ana"), str(tmp_path))

    assert Path(created).exists()


def test_generate_ficha_pptx_uses_normalized_filename(tmp_path: Path) -> None:
    employee = _employee(nome="JoÃ£o BÃ¡rbara")

    created = generator_ficha.generate_ficha_pptx(employee, str(tmp_path))

    assert Path(created).stem == normalize_filename(employee["nome"])


def test_generate_ficha_pptx_reports_single_progress_event(tmp_path: Path) -> None:
    messages: list[dict] = []

    created = generator_ficha.generate_ficha_pptx(
        _employee(nome="Ana"),
        str(tmp_path),
        callback=messages.append,
    )

    assert Path(created).exists()
    assert any(
        message.get("type") == "progress"
        and message.get("current") == 1
        and message.get("total") == 1
        for message in messages
    )
