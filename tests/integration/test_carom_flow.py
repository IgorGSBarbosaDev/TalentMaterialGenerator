from __future__ import annotations

from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation

from app.core.generator_carom import CaromConfig, generate_carom_pptx
from app.core.reader import load_standardized_carom_rows, read_spreadsheet


def _build_standardized_carom_spreadsheet(path: Path, total_rows: int) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.append(
        [
            "Matricula",
            "Nome",
            "Idade",
            "Cargo",
            "Formacao",
            "Nota 2025",
            "Potencial 2025",
            "CEO3",
            "CEO4",
        ]
    )
    for index in range(1, total_rows + 1):
        sheet.append(
            [
                str(100 + index),
                f"Colab {index}",
                str(20 + index),
                "Analista",
                "Engenharia",
                "4",
                "AP",
                f"CEO3 {index}",
                f"CEO4 {index}",
            ]
        )
    workbook.save(path)
    return path


def _build_talent_review_spreadsheet_without_ceos(path: Path, total_rows: int) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.append(
        [
            "Matricula",
            "Nome",
            "Idade",
            "Cargo",
            "Formacao",
            "Nota 2025",
            "Potencial 2025",
        ]
    )
    for index in range(1, total_rows + 1):
        sheet.append(
            [
                str(100 + index),
                f"Colab {index}",
                str(20 + index),
                "Analista",
                "Engenharia",
                "4",
                "AP",
            ]
        )
    workbook.save(path)
    return path


def _build_legacy_standardized_carom_spreadsheet(path: Path, total_rows: int) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.append(["Matricula", "Nome", "Idade", "Cargo"])
    for index in range(1, total_rows + 1):
        sheet.append(
            [
                str(100 + index),
                f"Colab {index}",
                str(20 + index),
                "Analista",
            ]
        )
    workbook.save(path)
    return path


def test_legacy_standardized_carom_flow_generates_regular_alias(tmp_path: Path) -> None:
    spreadsheet = _build_legacy_standardized_carom_spreadsheet(
        tmp_path / "carom-legacy.xlsx",
        3,
    )
    employees = load_standardized_carom_rows(read_spreadsheet(str(spreadsheet)))
    config: CaromConfig = {
        "preset_id": "regular",
        "titulo": "Carometro",
        "file_basename": "Carometro",
    }

    generated_files = generate_carom_pptx(employees, str(tmp_path), config)

    assert len(generated_files) == 1
    assert Path(generated_files[0]).exists()


def test_full_big_carom_flow_generates_single_file(tmp_path: Path) -> None:
    spreadsheet = _build_standardized_carom_spreadsheet(tmp_path / "carom.xlsx", 3)
    employees = load_standardized_carom_rows(read_spreadsheet(str(spreadsheet)))
    config: CaromConfig = {
        "preset_id": "big",
        "titulo": "Carometro",
        "file_basename": "Carometro",
    }

    generated_files = generate_carom_pptx(employees, str(tmp_path), config)

    assert len(generated_files) == 1
    assert Path(generated_files[0]).exists()


def test_big_carom_flow_creates_expected_slide_count(tmp_path: Path) -> None:
    spreadsheet = _build_standardized_carom_spreadsheet(tmp_path / "carom-big.xlsx", 9)
    employees = load_standardized_carom_rows(read_spreadsheet(str(spreadsheet)))
    config: CaromConfig = {
        "preset_id": "big",
        "titulo": "Carometro",
        "file_basename": "Carometro",
    }

    generated_files = generate_carom_pptx(employees, str(tmp_path), config)
    presentation = Presentation(generated_files[0])

    assert len(presentation.slides) == 2


def test_projeto_trainee_flow_uses_literal_body_text(tmp_path: Path) -> None:
    spreadsheet = _build_standardized_carom_spreadsheet(
        tmp_path / "carom-trainee.xlsx", 2
    )
    employees = load_standardized_carom_rows(read_spreadsheet(str(spreadsheet)))
    config: CaromConfig = {
        "preset_id": "projeto_trainee",
        "titulo": "Ignored",
        "file_basename": "Projeto_Trainee",
    }

    generated_files = generate_carom_pptx(employees, str(tmp_path), config)
    slide = Presentation(generated_files[0]).slides[0]

    assert slide.shapes[9].text.splitlines()[0] == "insira projeto trainee aqui"
    assert slide.shapes[13].text.splitlines()[0] == "insira projeto trainee aqui"


def test_talent_review_flow_keeps_template_text_without_ceo_fields(
    tmp_path: Path,
) -> None:
    spreadsheet = _build_talent_review_spreadsheet_without_ceos(
        tmp_path / "carom-tr.xlsx", 1
    )
    employees = load_standardized_carom_rows(read_spreadsheet(str(spreadsheet)))
    config: CaromConfig = {
        "preset_id": "talent_review",
        "titulo": "Ignored",
        "file_basename": "Talent_Review",
    }

    generated_files = generate_carom_pptx(employees, str(tmp_path), config)
    slide = Presentation(generated_files[0]).slides[0]
    text_box = slide.shapes[5]
    paragraphs = [paragraph.text for paragraph in text_box.text_frame.paragraphs]
    slide_text = "\n".join(
        child.text
        for shape in slide.shapes
        for child in (
            [shape] + list(shape.shapes) if hasattr(shape, "shapes") else [shape]
        )
        if hasattr(child, "text")
    )

    assert paragraphs[2] == "Sucessor Imediato"
    assert paragraphs[3] == "NomeCadeira"
    assert paragraphs[4] == "Em desenvolvimento"
    assert paragraphs[5] == "NomeCadeira"
    assert "CEO3" not in slide_text
    assert "CEO4" not in slide_text


def test_talent_review_flow_ignores_ceo_values_when_columns_are_present(
    tmp_path: Path,
) -> None:
    spreadsheet = _build_standardized_carom_spreadsheet(
        tmp_path / "carom-tr-ceos.xlsx", 12
    )
    employees = load_standardized_carom_rows(read_spreadsheet(str(spreadsheet)))
    config: CaromConfig = {
        "preset_id": "talent_review",
        "titulo": "Ignored",
        "file_basename": "Talent_Review",
    }

    generated_files = generate_carom_pptx(employees, str(tmp_path), config)
    slide = Presentation(generated_files[0]).slides[0]
    slide_text = "\n".join(
        child.text
        for shape in slide.shapes
        for child in (
            [shape] + list(shape.shapes) if hasattr(shape, "shapes") else [shape]
        )
        if hasattr(child, "text")
    )

    assert slide_text.count("Sucessor Imediato") == 12
    assert slide_text.count("Em desenvolvimento") == 12
    assert slide_text.count("NomeCadeira") == 24
    assert "CEO3" not in slide_text
    assert "CEO4" not in slide_text
